# -*- coding: utf-8 -*-

import sys

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

def my_paragraph(slide, loc, paragraphs, pt=10.5, bold=False):
	""" テキストボックスの挿入 """
	
	# テキストボックスの新規作成(サイズは適当)
	left = loc[0]
	top = loc[1]
	width = Cm(1.0)
	height = Cm(1.0)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	#tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
	
	# テキスト挿入
	max_w_len = 0
	for i, phrase in enumerate(paragraphs):
		if i==0:
			p = tf.paragraphs[0]
		else:
			p = tf.add_paragraph()
		
		p.text = phrase
		p.font.size = Pt(pt)
		p.font.bold = bold
		#p.font.color.rgb = RGBColor(255, 0, 0)
		
		# サイズ計算用に最大数を保持
		max_w_len = max(max_w_len, len(phrase))
	
	# 1文字のサイズ
	# pt 10.5 = 0.37x0.45cm
	# pt 14 = 0.5cmx0.6cm
	# マージン 0.25cmx2, 0.13cmx2
	
	# サイズの計算(width)
	w = {10.5:0.37, 14:0.5}
	txBox.width = Cm( w[pt]*max_w_len + 0.25*2 )
	
	# サイズの計算(height)
	h = {10.5:0.45, 14:0.6}
	txBox.height = Cm( h[pt]*len(paragraphs) + 0.13*2 )
	

if __name__ == "__main__":

	#
	# プレゼンテーションの新規作成
	#
	prs = Presentation()

	# スライドの挿入(タイトル)
	title_slide = prs.slides.add_slide( prs.slide_layouts[0] )

	# タイトルのテキストを記入
	title = title_slide.shapes.title
	title.text = 'python-pptx サンプルのタイトル'

	# サブタイトルのテキストを記入
	subtitle = title_slide.placeholders[1]
	subtitle.text = 'python-pptx サンプルのサブタイトル'
	
	# スライドのサイズ
	slide_h, slide_w = Emu(prs.slide_height), Emu(prs.slide_width)
	
	
	#
	# プロローグ
	#
	
	# スライドの挿入(ブランク)
	slide = prs.slides.add_slide( prs.slide_layouts[6] )
	
	# テキストボックスの挿入
	left = Cm(1.0)
	top = Cm(1.0)
	loc=[left, top]
	phrase  =['プロローグ']
	my_paragraph(slide, loc, phrase, pt=14, bold=True)
	

	

	# 画像の挿入
	#pic = slide.shapes.add_picture('サンプル.bmp', Cm(1.0), Cm(5.0))
	#pic = slide.shapes.add_picture('サンプル.bmp', Cm(1.0), Cm(10.0), height=Cm(5.0))
	#pic = slide.shapes.add_picture('序章-聖戦士誕生-プロローグ.gif', Cm(1.0), Cm(5.0))
	
	
	#
	# Phase1
	#
	# インデント(マージン)
	l0 = Cm(0.5)
	
	# スライドの挿入(ブランク)
	slide = prs.slides.add_slide( prs.slide_layouts[6] )
	
	# 表題の挿入(left, top, width, height)
	loc=[l0, Cm(1.0)]
	phrase  =['Phase1']
	my_paragraph(slide, loc, phrase, pt=14, bold=True)
	
	
	
	
	
	
	
	# スライドの挿入(ブランク)
	slide = prs.slides.add_slide( prs.slide_layouts[6] )
	
	# マップの挿入
	base_dir = '/storage/43E3-332A/others/test/work/ROM/SNES/ファイアーエムブレム 聖戦の系譜 (J)/'
	sub_dir = '00-序章-聖戦士誕生'
	filename = base_dir + sub_dir + '/序章-聖戦士誕生-フィールド-phase1(ユングヴィ城).png'
	pic = slide.shapes.add_picture(filename, Cm(0.0), Cm(0.0))
	pic.width = int(pic.width/3.5)
	pic.height = int(pic.height/3.5)
	pic.left =int( (slide_w-pic.width)/2 )
	#pic.top = int( (slide_h-pic.height)/2 )
	pic.top = Cm(1.0)
	
	# 表題の挿入(left, top, width, height)
	loc=[l0, Cm(1.0)]
	phrase  =['Phase1']
	my_paragraph(slide, loc, phrase, pt=14, bold=True)
	
	# 軍(左)
	loc=[l0, Cm(2.0)]
	phrase  =['シグルド軍(10)']
	my_paragraph(slide, loc, phrase, pt=10.5, bold=True)

	
	# 軍(左)の指揮官画像
	pic = slide.shapes.add_picture('temp/00.png', l0, Cm(2.6))
	w, h = pic.width, pic.height # 元画像のサイズを保持しておく
	pic.width = int(w*1.5)
	pic.height = int(h*1.5)
	
	# 指揮官以外の画像の挿入
	img_no =[
		[1, 2, 3],
		[4, 5],
		[6, 7, 8],
		[10]
	]
	texts = [
		'',
		'(2ターン目に加入)',
		'(3ターン目に加入)',
		'(ユングヴィ城制圧後に加入)',
	]
	wi, hi = int(w*0.75), int(h*0.75)
	for im, txt in zip(img_no, texts):
		# 注釈
		top =  pic.top + pic.height
		loc = [l0, top]
		if txt!='':
			phrase  =[txt]
			my_paragraph(slide, loc, phrase, pt=10.5, bold=False)

		
	
		# 画像
		if txt=='':
			top =  pic.top + pic.height + Cm(0.3)
		else:
			top =  pic.top + pic.height + Cm(0.6)
		for i, m in enumerate(im):
			left = l0 + (wi+Cm(0.3)) * i 
			pic = slide.shapes.add_picture('temp/{0:0>2}.png'.format(m), left, top, wi, hi)
	
	# 軍(右)の指揮官画像
	pic = slide.shapes.add_picture('temp/13.png', Cm(0.0), Cm(2.6))
	pic.width = int(w*1.5)
	pic.height = int(h*1.5)
	#l1 = slide_w - l0 - pic.width
	l1 = slide_w - Cm(4.0) - pic.width
	pic.left = l1
	
	# 軍(右)
	loc=[l1, Cm(2.0)]
	phrase  =['デマジオ軍(30)']
	my_paragraph(slide, loc, phrase, pt=10.5, bold=True)

	
	
	# 指揮官以外の画像の挿入
	img_no =[
		[14, 15, 16, 17, 18, 19, 22, 23],
		[24, 25, 26, 27, 28, 29, 30, 34],
		[35, 37, 38, 41, 42],
		[20, 21, 39],
		[31, 32, 33, 36, 40],
	]
	texts = [
		'バーバリアン(21)',
		'',
		'',
		'アーチャー(3)',
		'マウンテンシーフ(5)',
	]
	wi, hi = int(w*0.375), int(h*0.375)
	for im, txt in zip(img_no, texts):
		# 注釈
		top =  pic.top + pic.height
		
		loc = [l1, top]
		if txt!='':
			phrase  =[txt]
			my_paragraph(slide, loc, phrase, pt=10.5, bold=False)
			
	
		# 画像
		if txt=='':
			top =  pic.top + pic.height + Cm(0.15)
		else:
			top =  pic.top + pic.height + Cm(0.6)
		for i, m in enumerate(im):
			left = l1 + (wi+Cm(0.15)) * i 
			pic = slide.shapes.add_picture('temp/{0:0>2}.png'.format(m), left, top, wi, hi)
	
	
	
	# セリフ(村)
	# スライドの挿入(ブランク)
	slide = prs.slides.add_slide( prs.slide_layouts[6] )
	import pptx01 as phr
	
	# 表題の挿入(left, top, width, height)
	loc=[l0, Cm(1.0)]
	phrase  =['村解放']
	my_paragraph(slide, loc, phrase, pt=14, bold=True)

	
	l1 = int(slide_w / 2) + l0
	phrases = [
		[(l0, Cm(2.0)), 2, (l0+Cm(2.0), Cm(2.6)), phr.phrase1],
		[(l0, Cm(6.0)), 42, (l0+Cm(2.0), Cm(6.6)), phr.phrase2],
		[(l0, Cm(11.0)), 92, (l0+Cm(2.0), Cm(11.6)), phr.phrase3],
		
		
		[(l1, Cm(2.0)), 132, (l1+Cm(2.0), Cm(2.6)), phr.phrase4],
		[(l1, Cm(8.0)), 192, (l1+Cm(2.0), Cm(8.6)), phr.phrase5],
	]
	
	for i, (loc0, im, loc1, phrase) in enumerate(phrases):
		txt  =['村{}'.format(i+1)]
		my_paragraph(slide, loc0, txt, pt=10.5, bold=True)
		loc = (loc0[0], loc1[1])
		pic = slide.shapes.add_picture('temp2/{0:0>3}.png'.format(im), *loc)
		my_paragraph(slide, loc1, phrase, pt=10.5, bold=False)

	
	
	#
	
	# Phase2
	#
	# スライドの挿入(ブランク)
	slide = prs.slides.add_slide( prs.slide_layouts[6] )
	
	# マップの挿入
	base_dir = '/storage/43E3-332A/others/test/work/ROM/SNES/ファイアーエムブレム 聖戦の系譜 (J)/'
	sub_dir = '00-序章-聖戦士誕生'
	filename = base_dir + sub_dir + '/序章-聖戦士誕生-フィールド-phase2(エバンズ城).png'
	pic = slide.shapes.add_picture(filename, Cm(0.0), Cm(0.0))
	pic.width = int(pic.width/3.5)
	pic.height = int(pic.height/3.5)
	pic.left =int( (slide_w-pic.width)/2 )
	#pic.top = int( (slide_h-pic.height)/2 )
	pic.top = Cm(1.0)
	
	# インデント(マージン)
	l0 = Cm(0.5)
	
	# 表題の挿入(left, top, width, height)
	loc=[l0, Cm(1.0)]
	phrase  =['Phase2']
	my_paragraph(slide, loc, phrase, pt=14, bold=True)


	# 軍(左)
	loc=[l0, Cm(2.0)]
	phrase  =['シグルド軍(10)']
	my_paragraph(slide, loc, phrase, pt=10.5, bold=True)
	
	# 軍(左)の指揮官画像
	pic = slide.shapes.add_picture('temp/00.png', l0, Cm(2.6))
	w, h = pic.width, pic.height # 元画像のサイズを保持しておく
	pic.width = int(w*1.5)
	pic.height = int(h*1.5)
	
	# 指揮官以外の画像の挿入
	img_no =[
		[1, 2, 3, 4, 5],
		[6, 7, 8, 10],
	]
	texts = [
		'',
		'',
	]
	wi, hi = int(w*0.75), int(h*0.75)
	for im, txt in zip(img_no, texts):
		# 注釈
		top =  pic.top + pic.height
		if txt!='':
			loc=[l0, top]
			phrase = [txt]
			my_paragraph(slide, loc, phrase, pt=10.5, bold=True)
	
		# 画像
		if txt=='':
			top =  pic.top + pic.height + Cm(0.3)
		else:
			top =  pic.top + pic.height + Cm(0.6)
		for i, m in enumerate(im):
			left = l0 + (wi+Cm(0.3)) * i 
			pic = slide.shapes.add_picture('temp/{0:0>2}.png'.format(m), left, top, wi, hi)
	
	# 軍(右)の指揮官画像
	pic = slide.shapes.add_picture('temp/43.png', Cm(0.0), Cm(2.6))
	pic.width = int(w*1.5)
	pic.height = int(h*1.5)
	#l1 = slide_w - l0 - pic.width
	l1 = slide_w - Cm(4.0) - pic.width
	pic.left = l1
	
	# 軍(右)
	loc=[l1, Cm(2.0)]
	phrase  =['ゲラルド軍(16)']
	my_paragraph(slide, loc, phrase, pt=10.5, bold=True)
	
	
	# 指揮官以外の画像の挿入
	img_no =[
		[49, 50, 51, 52, 53, 54, 55, 56],
		[57, 58],
		[44, 45, 46, 47, 48],
	]
	texts = [
		'バーバリアン(10)',
		'',
		'アーチャー(5)',
	]
	wi, hi = int(w*0.375), int(h*0.375)
	for im, txt in zip(img_no, texts):
		# 注釈
		top =  pic.top + pic.height
		
		if txt!='':
			loc=[l1, top]
			phrase = [txt]
			my_paragraph(slide, loc, phrase, pt=10.5, bold=True)
	
		# 画像
		if txt=='':
			top =  pic.top + pic.height + Cm(0.15)
		else:
			top =  pic.top + pic.height + Cm(0.6)
		for i, m in enumerate(im):
			left = l1 + (wi+Cm(0.15)) * i 
			pic = slide.shapes.add_picture('temp/{0:0>2}.png'.format(m), left, top, wi, hi)
	
	


	# プレゼンテーションの保存
	prs.save('サンプル.pptx')
