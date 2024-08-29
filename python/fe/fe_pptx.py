# -*- coding: utf-8 -*-

import sys
import glob

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor


class fe_pptx(object):
	""" プレゼンテーション作成 """
	TEMP_DIR = 'temp'
	def __init__(self, title):
		""" 初期化 """
		# プレゼンテーションの作成
		self.prs = Presentation()
		
		# 表題
		self.title = title
	
	def save(self):
		""" プレゼンテーションの保存 """
		self.prs.save( self.title + '.pptx' )
	
	def add_text(self, slide, loc, paragraphs, pt=10.5, bold=False):
		""" テキストボックスの挿入 """
		# テキストボックスの新規作成(サイズは適当)
		left = loc[0]
		top = loc[1]
		width = Cm(1.0)
		height = Cm(1.0)
		
		txBox = slide.shapes.add_textbox(left, top, width, height)
		tf = txBox.text_frame
		#tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
	
		# テキスト挿入
		max_w_len = 0
		for i, phrase in enumerate(paragraphs):
			if i==0:
				p = tf.paragraphs[0]
			else:
				p = tf.add_paragraph()
		
			p.text = phrase
			p.font.size = Pt(pt)
			p.font.bold = bold
			#p.font.color.rgb = RGBColor(255, 0, 0)
		
			# サイズ計算用に最大数を保持
			max_w_len = max(max_w_len, len(phrase))
	
		# 1文字のサイズ
		# pt 10.5 = 0.37x0.45cm
		# pt 14 = 0.5cmx0.6cm
		# マージン 0.25cmx2, 0.13cmx2
	
		# サイズの計算(width)
		w = {10.5:0.37, 14:0.5}
		txBox.width = Cm( w[pt]*max_w_len + 0.25*2 )
	
		# サイズの計算(height)
		h = {10.5:0.45, 14:0.6}
		txBox.height = Cm( h[pt]*len(paragraphs) + 0.13*2 )
		
		return txBox.width, txBox.height
	
	def add_slide_line(self, text):
		""" スライド中央に線を引いて追加 """
		slide = self.prs.slides.add_slide( self.prs.slide_layouts[6] )
		self.add_text(slide, [Cm(1.0), Cm(1.0)], [text], pt=14, bold=True)
		slide.shapes.add_shape(
			MSO_SHAPE_TYPE.LINE,
			int(self.prs.slide_width /2),
			Cm(2.0),
			0,
			self.prs.slide_height - Cm(2.0)*2
			)

		return slide
	
	def make_title(self, **kw):
		""" タイトルスライドの作成 """
		# スライドの挿入(タイトル)
		slide = self.prs.slides.add_slide( self.prs.slide_layouts[0] )

		# タイトルのテキストを記入
		title = slide.shapes.title
		title.text = 'ファイヤーエムブレム'
		title.text_frame.add_paragraph().text = '聖戦の系譜'

		# サブタイトルのテキストを記入
		sub = slide.placeholders[1]
		sub.text = self.title
		
		if 'subtitle' in kw.keys():
			# 空行
			sub.text_frame.add_paragraph()
			
			for st in kw['subtitle']:
				p = sub.text_frame.add_paragraph()
				p.text = st
				p.font.size = Pt(18)
	
	def make_ivent(self, **kw):
		""" イベントシーンの作成 """
		# スライド作成(中央にライン)
		slide = self.add_slide_line(kw['title'])
		
		# 画像とテキストの貼り付け
		l0 = Cm(0.5)
		t0 = Cm(2.0)
		l1, t1 = l0, t0
		for i, (im, a, paragraphs) in enumerate(kw['scene']):
			# 改ページ
			if a==-1:
				t1 = t0
				if l1==l0:
					l1 = l0 + int(self.prs.slide_width/2)
				else:
					l1 = l0
					slide = self.add_slide_line(kw['title'])
				
				continue

			# 画像
			if im=='':
				# 画像がない時はテキストのみ
				pic_w = -Cm(0.3)
				pic_h = 0
			else:
				filename = self.TEMP_DIR + '/{0}.png'.format(im)
				pic = slide.shapes.add_picture(filename, l1, t1)
				pic.width = int(pic.width * a)
				pic.height = int(pic.height * a)
				pic_w = pic.width
				pic_h = pic.height
			
			# テキスト
			if paragraphs==[]:
				# テキストがない時は中央に
				pic.left = (l1 - l0) + int(self.prs.slide_width/4 - pic_w/2)
				w, h = 0, 0
			else:
				# テキストがある時は貼り付け
				l1_ = l1 + pic_w + Cm(0.3)
				w, h = self.add_text(slide, [l1_, t1], paragraphs)
			
			if pic_h > h:
				t1 += pic_h +Cm(0.5)
			else:
				t1 +=  h + Cm(0.5)
			
			# はみ出す時の処理はTBD->改ページで対応
			if t1 + pic_h >= self.prs.slide_height:
				t1 = t0
				if l1==l0:
					l1 = l0 + int(self.prs.slide_width/2)
				else:
					l1 = l0
					if i!=len(kw['scene'])-1:
						# 終わりでなければスライド作成
						slide = self.add_slide_line(kw['title'])
	
	def make_prologue(self, **kw):
		""" プロローグスライドの作成 """
		# 全体マップ
		# スライドの挿入(ブランク)
		slide = self.prs.slides.add_slide( self.prs.slide_layouts[6] )
		self.add_text(slide, [Cm(1.0), Cm(1.0)], ['全体マップ'], pt=14, bold=True)
		
		# 全体マップは中央に配置
		pic = slide.shapes.add_picture(kw['prologue_map'], Cm(1.0), Cm(2.0))
		pic.left = int( (self.prs.slide_width - pic.width)/2 )
		pic.top = int( (self.prs.slide_height - pic.height)/2 )
		
		if pic.left<Cm(1.0) or pic.top<Cm(2.0):
			l = pic.left - Cm(1.0)
			t = pic.top - Cm(2.0)
			
			if l>t:
				a = pic.width / pic.height
				pic.height = self.prs.slide_height - Cm(2.0)*2
				pic.width = int(a * pic.height)
				pic.left = int( (self.prs.slide_width - pic.width)/2 )
				pic.top = Cm(2.0)
			else:
				a = pic.height / pic.width
				pic.width = self.prs.slide_width - Cm(1.0)*2
				pic.height = int(a * pic.width)
				pic.left = Cm(1.0)
				pic.top = int( (self.prs.slide_height - pic.height)/2 )
		
		# イベントシーン
		self.make_ivent(**kw)
				
		# 本当はアニメを貼りたい
		#slide.shapes.add_movie('./test.gif', Cm(1.0), Cm(1.0), Cm(10.0), Cm(15.0), mime_type='image/gif')
		
		# フィールドマップ
		slide = self.prs.slides.add_slide( self.prs.slide_layouts[6] )
		self.add_text(slide, [Cm(1.0), Cm(1.0)], ['フィールドマップ'], pt=14, bold=True)
	
		# フィールドマップは中央に配置
		pic = slide.shapes.add_picture(kw['field_map'], Cm(1.0), Cm(2.0))
		pic.left = int( (self.prs.slide_width - pic.width)/2 )
		pic.top = int( (self.prs.slide_height - pic.height)/2 )
		
		if pic.left<Cm(1.0) or pic.top<Cm(2.0):
			l = pic.left - Cm(1.0)
			t = pic.top - Cm(2.0)
			
			if l>t:
				a = pic.width / pic.height
				pic.height = self.prs.slide_height - Cm(2.0)*2
				pic.width = int(a * pic.height)
				pic.left = int( (self.prs.slide_width - pic.width)/2 )
				pic.top = Cm(2.0)
			else:
				a = pic.height / pic.width
				pic.width = self.prs.slide_width - Cm(1.0)*2
				pic.height = int(a * pic.width)
				pic.left = Cm(1.0)
				pic.top = int( (self.prs.slide_height - pic.height)/2 )
		
	def make_organization(self, **kw):
		""" 編成 """
		# スライドの挿入(ブランク)
		slide = self.prs.slides.add_slide( self.prs.slide_layouts[6] )
		self.add_text(slide, [Cm(1.0), Cm(1.0)], [kw['title']], pt=14, bold=True)
		
		# マップの貼り付け
		pic = slide.shapes.add_picture(kw['field_map'], Cm(0.0), Cm(0.0))
		pic.width = int(pic.width/3.5)
		pic.height = int(pic.height/3.5)
		pic.left =int( (self.prs.slide_width-pic.width)/2 )
		#pic.top = int( (self.prs.slide_height-pic.height)/2 )
		pic.top = Cm(2.3)
		
		# 軍(左)
		l0 = Cm(0.5)
		loc = [l0, Cm(2.0)]
		phrase = [ kw['organization']['左']['名'] ]
		self.add_text(slide, loc, phrase, pt=10.5, bold=True)
		
		# 軍(左)の指揮官画像
		im = kw['organization']['左']['指揮官']
		filename = self.TEMP_DIR + '/{0}.png'.format(im)
		pic = slide.shapes.add_picture(filename, l0, Cm(2.6))
		w, h = pic.width, pic.height # 元画像のサイズを保持しておく
		pic.width = int(w*1.5)
		pic.height = int(h*1.5)
		
		# 編成(左)
		if '倍率' in kw['organization']['左']:
			alpha = kw['organization']['左']['倍率']
		else:
			alpha = 0.75
		wi, hi = int(w*alpha), int(h*alpha)
		for txt, imgs in kw['organization']['左']['編成']:
			top =  pic.top + pic.height
			loc = [l0, top]
			
			# 注釈
			if txt!='':
				phrase  =[txt]
				self.add_text(slide, loc, phrase, pt=10.5, bold=False)
				
			# 画像
			if txt=='':
				top =  pic.top + pic.height + Cm(0.15)
			else:
				top =  pic.top + pic.height + Cm(0.6)
				
			for i, im in enumerate(imgs):
				left = l0 + (wi+Cm(0.15)) * i 
				filename = self.TEMP_DIR + '/{0:0>2}.png'.format(im)
				pic = slide.shapes.add_picture(filename, left, top, wi, hi)
				
		# 軍(右)の指揮官画像
		im = kw['organization']['右']['指揮官']
		filename = self.TEMP_DIR + '/{0}.png'.format(im)
		pic = slide.shapes.add_picture(filename, l0, Cm(2.6))
		pic.width = int(w*1.5)
		pic.height = int(h*1.5)
		#l1 = self.prs.slide_width - l0 - pic.width
		l1 = self.prs.slide_width - Cm(4.0) - pic.width
		pic.left = l1
		
		# 軍(右)
		loc=[l1, Cm(2.0)]
		phrase = [ kw['organization']['右']['名'] ]
		self.add_text(slide, loc, phrase, pt=10.5, bold=True)
		
		# 編成(右)
		wi, hi = int(w*0.375), int(h*0.375)
		for txt, imgs in kw['organization']['右']['編成']:
			top =  pic.top + pic.height
			loc = [l1, top]
			
			# 注釈
			if txt!='':
				phrase  =[txt]
				self.add_text(slide, loc, phrase, pt=10.5, bold=False)
				
			# 画像
			if txt=='':
				top =  pic.top + pic.height + Cm(0.15)
			else:
				top =  pic.top + pic.height + Cm(0.6)
				
			for i, im in enumerate(imgs):
				left = l1 + (wi+Cm(0.15)) * i 
				filename = self.TEMP_DIR + '/{0:0>2}.png'.format(im)
				pic = slide.shapes.add_picture(filename, left, top, wi, hi)

if __name__ == "__main__":
	pass