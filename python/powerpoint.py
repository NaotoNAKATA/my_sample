# -*- coding: utf-8 -*-

import sys

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

if __name__ == "__main__":

	#
	# プレゼンテーションの新規作成
	#
	prs = Presentation()

	# スライドの挿入(タイトル)
	title_slide = prs.slides.add_slide( prs.slide_layouts[0] )

	# タイトルのテキストを記入
	title = title_slide.shapes.title
	title.text = 'python-pptx サンプルのタイトル'

	# サブタイトルのテキストを記入
	subtitle = title_slide.placeholders[1]
	subtitle.text = 'python-pptx サンプルのサブタイトル'

	# スライドの挿入(ブランク)
	slide = prs.slides.add_slide( prs.slide_layouts[6] )

	# テキストボックスの挿入
	left = Cm(1.0)
	top = Cm(1.0)
	width = Cm(10.0)
	height = Cm(1.0)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	#tf.text = 'テキストボックス'
	p = tf.paragraphs[0]
	p.text = 'テキストボックス'
	p.font.size = Pt(10.5)
	p.font.bold = True
	p.font.color.rgb = RGBColor(255, 0, 0)

	p = tf.add_paragraph()
	p.text = '追加のパラグラフ'
	p.font.size = Pt(20)
	p.font.bold = True
	p.font.color.rgb = RGBColor(0, 255, 0)

	# 画像の挿入
	pic = slide.shapes.add_picture('サンプル.bmp', Cm(1.0), Cm(5.0))
	pic = slide.shapes.add_picture('サンプル.bmp', Cm(1.0), Cm(10.0), height=Cm(5.0))

	# 図形の挿入
	slide = prs.slides.add_slide( prs.slide_layouts[6] )
	left = Cm(1.0)
	top = Cm(1.0)
	width = Cm(4.0)
	height = Cm(1.0)

	for n in range(1, 6):
		if n==1:
			shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
		else:
			shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
		shape.text = '図形{0}'.format(n)
		left = left + width
	
	# 表の挿入
	left = Cm(1.0)
	top = Cm(3.0)
	width = Cm(4.0)
	height = Cm(1.0)
	rows = 2
	cols = 2
	table = slide.shapes.add_table(rows, cols, left, top, width, height).table

	table.columns[0].width = Cm(4.0)
	table.columns[1].width = Cm(4.0)
	for col in range(cols):
		for row in range(rows):
			table.cell(row, col).text = 'テスト{0}{1}'.format(row, col)
			p = table.cell(row, col).text_frame.paragraphs[0]
			p.font.size = Pt(8)
			p.font.bold = True
			p.alignment = PP_ALIGN.CENTER

	# プレゼンテーションの保存
	prs.save('サンプル.pptx')

